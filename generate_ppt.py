import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

# Define slides
slides_data = [
    {
        "title": "AI Integrated Eye Diagnosis and Doctor Consultation Management System",
        "content": "Department of Computer Science and Engineering\n\nTeam Members:\n[Student 1 Name] - [Roll Number]\n[Student 2 Name] - [Roll Number]\n[Student 3 Name] - [Roll Number]\n\nGuided By: [Faculty Name] - [Designation]",
        "is_title": True
    },
    {
        "title": "Objective",
        "content": [
            "To develop an automated, AI-driven platform for early detection of eye diseases (such as Diabetic Retinopathy, Glaucoma, Cataract).",
            "To provide a seamless, integrated portal for patients to consult with certified ophthalmologists online.",
            "To reduce diagnosis time and make specialized eye care accessible to remote and underserved populations.",
            "To maintain secure and centralized digital medical records for patients and doctors."
        ]
    },
    {
        "title": "Abstract",
        "content": [
            "The project introduces an AI-integrated telehealth system tailored for ophthalmology.",
            "Utilizing advanced Deep Learning and Computer Vision, the system analyzes retinal images to accurately identify common eye conditions.",
            "Patients can instantly view predictive results and subsequently book appointments or initiate remote consultations with medical professionals via the platform.",
            "The dual-approach combining automated AI screening with expert human validation ensures highly reliable medical diagnoses.",
            "This system significantly lowers the barrier to timely, cost-effective eye care through a user-friendly web interface."
        ]
    },
    {
        "title": "Existing System",
        "content": [
            "Current diagnosis primarily relies on manual examination of retinal fundus images by highly trained ophthalmologists.",
            "Patients must physically visit clinics and specialized eye hospitals to get screened for diseases.",
            "Rural or remote clinics often lack the specialized equipment and available specialists.",
            "Medical records and imaging history are frequently siloed or locally maintained."
        ]
    },
    {
        "title": "Drawbacks of Existing System",
        "content": [
            "Time-consuming: Requires manual observation of complex retinal scans.",
            "Human Error: Prone to fatigue and inter-observer variability in diagnosis.",
            "Resource Shortage: Severe shortage of expert ophthalmologists, especially in rural areas.",
            "Delayed Diagnosis: Can lead to irreversible vision loss before appropriate treatment.",
            "High Costs: Associated with frequent in-person hospital visits.",
            "Data Silos: Lack of centralized, easily accessible digital health records."
        ]
    },
    {
        "title": "Proposed System",
        "content": [
            "An integrated web application combining AI-powered computer vision with a telehealth portal.",
            "Users securely upload their retinal scans directly to the platform.",
            "Deep learning backend (CNNs) processes imagery to predict eye diseases.",
            "Connects patients seamlessly with registered doctors to share AI-generated reports.",
            "Features complete appointment management and secure digital medical records."
        ]
    },
    {
        "title": "Advantages of Proposed System",
        "content": [
            "Early and Rapid Detection: Instantaneous AI screening for faster intervention.",
            "High Accessibility: Access expert-level screening through a simple web browser.",
            "Reduced Workload: AI pre-screens and flags critical cases for specialists.",
            "Convenient Telemedicine: Built-in doctor consultation eliminates unnecessary visits.",
            "Centralized Health Records: Secure tracking of history, reports, and prescriptions."
        ]
    },
    {
        "title": "Problem Statement",
        "content": [
            "To design and implement an accessible, efficient, and accurate web-based system that addresses the critical problems of delayed eye disease diagnosis and the shortage of available specialists by leveraging Artificial Intelligence for automated retinal screening and providing seamlessly integrated remote doctor consultations."
        ]
    },
    {
        "title": "Literature Survey",
        "content": [
            "1. Gulshan et al. - Deep Learning for DR Detection: High sensitivity, but computationally heavy.",
            "2. Kumar et al. - Tele-ophthalmology in rural areas: Improved access, but slow manual grading.",
            "3. Ting et al. - AI in Ophthalmology: Comprehensive robustness, lacking clinical integration.",
            "4. Zhang et al. - Cloud Medical Image Processing: Scalable, but faces data privacy concerns."
        ]
    },
    {
        "title": "Technologies Used",
        "content": [
            "Operating System: Windows / Linux / macOS",
            "Programming Language: Python, JavaScript",
            "AI/ML Framework: TensorFlow, Keras, Scikit-learn",
            "Computer Vision Tools: OpenCV, PIL",
            "Backend: Flask (Python)",
            "Frontend: HTML5, CSS3, JavaScript, Bootstrap",
            "Database: SQLite / PostgreSQL",
            "Deployment Platform: Render / AWS / Vercel"
        ]
    },
    {
        "title": "System Architecture",
        "content": [
            "Step 1: User accesses Frontend to upload an image.",
            "Step 2: Frontend securely sends image to Flask Backend via HTTP.",
            "Step 3: Flask routes data to TensorFlow AI Model for inference.",
            "Step 4: AI Model returns predicted disease & confidence score.",
            "Step 5: Database saves Diagnosis Report, UI updates.",
            "Step 6: User navigates to Doctor Consultation Module to book an appointment."
        ]
    },
    {
        "title": "Identified Modules",
        "content": [
            "User Management Module",
            "Image Upload Module",
            "AI Eye Disease Detection Module",
            "Doctor Consultation Module",
            "Appointment Management Module",
            "Diagnosis Report Module",
            "Admin Management Module"
        ]
    },
    {
        "title": "Module Description",
        "content": [
            "User/Doctor Management: Registration, authentication, dashboards.",
            "Image Upload: Secure high-res eye scan uploads.",
            "AI Detection: ML engine for predictive classification.",
            "Doctor Consultation: Directory and profiling of specialists.",
            "Appointment Management: Scheduling and managing online visits.",
            "Diagnosis Report: comprehensive medical reports.",
            "Admin Management: Oversees roles, verifies credentials."
        ]
    },
    {
        "title": "Sample Code for AI Model",
        "content": [
            "import tensorflow as tf",
            "from tensorflow.keras.preprocessing import image",
            "import numpy as np",
            "",
            "model = tf.keras.models.load_model('eye_disease_model.h5')",
            "",
            "def predict_disease(img_path):",
            "    img = image.load_img(img_path, target_size=(224, 224))",
            "    img_array = image.img_to_array(img) / 255.0",
            "    predictions = model.predict(np.expand_dims(img_array, axis=0))",
            "    classes = ['Cataract', 'DR', 'Glaucoma', 'Normal']",
            "    return classes[np.argmax(predictions[0])]"
        ]
    },
    {
        "title": "Conclusion",
        "content": [
            "Formulated an robust solution for early eye disease detection.",
            "Integrated state-of-the-art deep learning models to reduce manual screening dependency.",
            "Bridged the healthcare gap with a functional doctor consultation portal.",
            "Delivered a secure centralized medical system improving patient convenience.",
            "Established a strong foundation for future telehealth expansions."
        ]
    },
    {
        "title": "Future Enhancements",
        "content": [
            "Mobile app integration for Android and iOS devices.",
            "Real-time telemedicine featuring built-in secure video conferencing.",
            "Cloud-based medical records unified with national health registry APIs.",
            "Advanced deep learning models capable of localized lesion segmentation.",
            "IoT integration directly with smart fundus cameras at rural medical camps."
        ]
    },
    {
        "title": "References",
        "content": [
            "Gulshan, V., et al. (2016). Deep Learning Algorithm for Detection of Diabetic Retinopathy.",
            "Ting, D. S. W., et al. (2017). Deep Learning System for Diabetic Retinopathy.",
            "Kumar, S., et al. (2020). Tele-ophthalmology: A Need of the Hour.",
            "TensorFlow Documentation: https://www.tensorflow.org/",
            "Flask Web Development: https://flask.palletsprojects.com/"
        ]
    }
]

title_slide_layout = prs.slide_layouts[0]
bullet_slide_layout = prs.slide_layouts[1]

for slide_data in slides_data:
    if slide_data.get("is_title"):
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = slide_data["title"]
        subtitle.text = slide_data["content"]
    else:
        slide = prs.slides.add_slide(bullet_slide_layout)
        title = slide.shapes.title
        body_shape = slide.placeholders[1]
        title.text = slide_data["title"]
        
        tf = body_shape.text_frame
        content = slide_data["content"]
        
        for i, point in enumerate(content):
            if i == 0:
                p = tf.paragraphs[0]
                p.text = point
            else:
                p = tf.add_paragraph()
                p.text = point
                
            # For sample code, make it smaller and monospace
            if slide_data["title"] == "Sample Code for AI Model":
                p.font.size = Pt(14)
                p.font.name = "Courier New"
                    
prs.save('Project_Presentation.pptx')
print("Successfully created Project_Presentation.pptx!")
